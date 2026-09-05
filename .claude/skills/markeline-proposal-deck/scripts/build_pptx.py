#!/usr/bin/env python3
"""Build an editable .pptx (native shapes + text, no raster slides) from an outline JSON.

    python3 scripts/build_pptx.py outline.json --out deck/deck.pptx

Opens in PowerPoint and imports into Google Slides with every text box editable.
Geometry is the same 1920x1080 px grid as the HTML template (1 px = 1/144 in), sizes come
from markeline-design-system/references/reference-type-scale.md, colors from tokens.py.
Fonts: by default every run is "Noto Sans JP" (bold where the reference uses Black) so the
file renders correctly the moment it is imported into Google Slides. `--weight black` asks
for the "Noto Sans JP Black" / "Noto Serif JP Black" faces instead (PowerPoint / Keynote
with the fonts installed).
"""
from __future__ import annotations

import argparse
import html as _html
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
DS = HERE.parents[1] / "markeline-design-system"
sys.path.insert(0, str(DS / "tokens"))
from tokens import HEX  # noqa: E402

ICON_DIR = DS / "assets" / "icons"
LOGO_COLOR = DS / "assets" / "markeline_logo_color.png"
LOGO_WHITE = DS / "assets" / "markeline_logo_white.png"

SANS_BLACK = "Noto Sans JP Black"
SANS = "Noto Sans JP"
SANS_MED = "Noto Sans JP Medium"
SERIF_BLACK = "Noto Serif JP Black"


# ---------- units / colors ----------
def px(v: float) -> Emu:
    return Inches(v / 144)


def col(name: str) -> RGBColor:
    h = HEX[name].lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def plain(s: str | None) -> str:
    return re.sub(r"\*\*|==|\\n|\n", "", s or "")


def fit(text, base, max_chars, floor, width):
    n = len(plain(text))
    return base if n <= max_chars else max(floor, int(width / n))


def nlines(content, size, width) -> int:
    """Estimate wrapped line count (CJK glyph ≈ 0.97 em) so boxes can be stacked without overlap."""
    s = content if isinstance(content, str) else "\n".join(content)
    s = plain(s.replace("\\n", "\n"))
    cpl = max(1, int(width / (size * 0.97)))
    return sum(max(1, -(-len(ln) // cpl)) for ln in s.split("\n")) or 1


# ---------- primitives ----------
def rect(slide, x, y, w, h, fill=None, line=None, radius=None, line_w=2):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    if radius:
        shp.adjustments[0] = min(0.5, radius / min(w, h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = col(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = col(line); shp.line.width = px(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def circle(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(d), px(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = col(fill)
    if line:
        shp.line.color.rgb = col(line); shp.line.width = px(3)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def polygon(slide, points, fill):
    x0 = min(p[0] for p in points); y0 = min(p[1] for p in points)
    fb = slide.shapes.build_freeform(px(points[0][0]), px(points[0][1]), scale=1.0)
    fb.add_line_segments([(px(x), px(y)) for x, y in points[1:]], close=True)
    shp = fb.convert_to_shape()
    shp.fill.solid(); shp.fill.fore_color.rgb = col(fill); shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def icon(slide, name, x, y, size, color="teal"):
    p = ICON_DIR / f"{name}_{color}.png"
    if not p.exists():
        p = ICON_DIR / f"target_{color}.png"
    return slide.shapes.add_picture(str(p), px(x), px(y), px(size), px(size))


BASE_DIR = Path(".")  # outline-relative `image` paths resolve against this


def picture(slide, path, x, y, w=None, h=None):
    p = Path(path)
    if not p.is_absolute() and not p.exists():
        p = BASE_DIR / p
    if not p.exists():
        print(f"warning: image not found: {path}", file=sys.stderr)
        return None
    return slide.shapes.add_picture(str(p), px(x), px(y), px(w) if w else None, px(h) if h else None)


BLACK_FACE = False  # --weight black: request the "… Black" faces (needs the font installed)


def _set_font(run, name, size_px, color, bold=False):
    if not BLACK_FACE:
        if name in (SANS_BLACK, SANS_MED):
            name, bold = SANS, True
        elif name == SERIF_BLACK:
            name, bold = "Noto Serif JP", True
    f = run.font
    f.name = name; f.size = Pt(size_px * 0.375); f.bold = bold; f.color.rgb = col(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def text(slide, x, y, w, h, content, size, color="ink", font=SANS, bold=False, align="l", valign="t",
         em_color="yellow-700", line=1.3, wrap=True, marker=None):
    """content: str with **em** / ==marker== / \\n markup, or list of lines."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    # wrap=False is only a hint: LibreOffice / Google Slides render wrap="none" boxes
    # grown around their centre, which shifts left-aligned text. Keep word wrap on and
    # rely on fit() / box widths to keep single-line content on one line.
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[valign]
    s = content if isinstance(content, str) else "\n".join(content)
    s = s.replace("\\n", "\n")
    first = True
    for ln in s.split("\n"):
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        para.line_spacing = line
        # tokens: **em**, ==marker==
        pos = 0
        for m in re.finditer(r"\*\*(.+?)\*\*|==(.+?)==", ln):
            if m.start() > pos:
                _set_font(para.add_run(), font, size, color, bold).text if False else None
                r = para.add_run(); r.text = ln[pos:m.start()]; _set_font(r, font, size, color, bold)
            r = para.add_run()
            if m.group(1) is not None:
                r.text = m.group(1); _set_font(r, font, size, em_color, bold)
            else:
                r.text = m.group(2); _set_font(r, font, size, color, bold)
                # highlight (yellow-100) via a:highlight
                rPr = r._r.get_or_add_rPr(); hl = rPr.makeelement(qn("a:highlight"), {})
                clr = hl.makeelement(qn("a:srgbClr"), {"val": HEX["yellow-100"].lstrip("#")}); hl.append(clr); rPr.append(hl)
            pos = m.end()
        if pos < len(ln) or not ln:
            r = para.add_run(); r.text = ln[pos:]; _set_font(r, font, size, color, bold)
    return tb


# ---------- chrome ----------
def logobox(slide, x=40, y=0, w=210, h=150, radius=8):
    rect(slide, x, y, w, h, fill="paper", radius=radius)
    picture(slide, LOGO_COLOR, x + (w - 140) / 2, y + (h - 100) / 2, 140, 100)


def titlebar(slide, sl):
    title = sl.get("title", "")
    if sl.get("sub"):
        h = 200
        polygon(slide, [(0, 0), (1920, 0), (1860, h), (0, h)], "teal-700")
        logobox(slide, 40, 0, 210, h)
        size = fit(title, 74, 20, 56, 1480)
        text(slide, 290, 18, 1520, 100, title, size, "paper", SANS_BLACK, valign="m", em_color="brand-yellow")
        text(slide, 290, 120, 1520, 60, sl["sub"], 40, "paper", SANS, bold=True, valign="m")
        return h
    h = 150
    polygon(slide, [(0, 0), (1920, 0), (1860, h), (0, h)], "teal-700")
    logobox(slide)
    size = fit(title, 88, 17, 60, 1540)
    text(slide, 290, 0, 1540, h, title, size, "paper", SANS_BLACK, valign="m", em_color="brand-yellow")
    return h


def conclusion(slide, c):
    if not c:
        return
    if isinstance(c, str):
        c = {"text": c}
    teal = c.get("tone") == "teal"
    size = fit(c.get("text"), 56, 28, 40, 1580)
    y = 930
    rect(slide, 0, y, 1920, 150, fill="teal-700" if teal else "brand-yellow")
    circle(slide, 72, y + 39, 72, "brand-yellow" if teal else "teal-700")
    icon(slide, "check", 72 + 14, y + 39 + 14, 44, "ink" if teal else "white")
    text(slide, 172, y, 1680, 150, c.get("text", ""), size, "paper" if teal else "ink", SANS_BLACK, valign="m",
         em_color="brand-yellow" if teal else "yellow-700")


def pageno(slide, page, total):
    text(slide, 1700, 1040, 180, 30, f"{page:02d} / {total:02d}", 22, "sub", SANS, align="r")


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    rect(slide, 0, 0, 1920, 1080, fill=color)


# ---------- slide types ----------
def s_cover(prs, sl, o, page, total):
    s = blank(prs); bg(s, "brand-teal")
    logobox(s, 60, 48, 210, 136)
    title = sl.get("title") or f'{o["client"]}{o.get("honorific", "様")} ご提案'
    if sl.get("image"):
        # illustration on the right, copy left-aligned (mirrors .cover.withimg in the template)
        x, w, al = 120, 1080, "l"
        picture(s, sl["image"], 1240, 260, 560, 560)
    else:
        x, w, al = 160, 1600, "c"
    size = fit(title, 126, 8 if sl.get("image") else 12, 88, w)
    if sl.get("label"):
        text(s, x, 330, w, 50, sl["label"], 32, "teal-100", SANS, bold=True, align=al)
    text(s, x, 390, w, 170, title, size, "paper", SANS_BLACK, align=al, valign="m")
    rect(s, x if al == "l" else 560, 575, 800, 4, fill="brand-yellow")
    text(s, x, 610, w, 80, f"〜 {sl.get('subtitle', '')} 〜", 48, "paper", SANS, bold=True, align=al, valign="m")
    text(s, 1200, 990, 648, 50, o.get("company", "株式会社MARKELINE"), 36, "paper", SANS, bold=True, align="r")


def s_goals(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    items = sl.get("items", [])[:3]
    h = 200; gap = 40
    y0 = top + ((1080 - top) - (len(items) * h + (len(items) - 1) * gap)) / 2
    for i, t in enumerate(items):
        y = y0 + i * (h + gap)
        rect(s, 72, y, 1776, h, fill="bg", radius=8)
        rect(s, 72, y, 12, h, fill="brand-yellow")
        circle(s, 132, y + (h - 150) / 2, 150, "teal-700")
        text(s, 132, y + (h - 150) / 2, 150, 150, str(i + 1), 110, "paper", SANS_BLACK, align="c", valign="m")
        text(s, 320, y, 1500, h, t, 64, "ink", SANS_BLACK, valign="m", line=1.35)
    pageno(s, page, total)


def s_issues(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    items = sl.get("items", [])[:4]
    cols = 3 if len(items) == 3 else 2
    rows = (len(items) + cols - 1) // cols
    gx = 24; gy = 24; x0 = 72; y0 = top + 32; W = 1776; H = 1080 - y0 - 32
    cw = (W - gx * (cols - 1)) / cols; ch = (H - gy * (rows - 1)) / rows
    for i, it in enumerate(items):
        x = x0 + (i % cols) * (cw + gx); y = y0 + (i // cols) * (ch + gy)
        rect(s, x, y, cw, ch, fill="bg", line="line", radius=8, line_w=1)
        icon(s, it.get("icon", "target"), x + 36, y + (ch - 140) / 2, 140, "teal")
        circle(s, x + 204, y + ch / 2 - 96, 90, "brand-yellow")
        text(s, x + 204, y + ch / 2 - 96, 90, 90, str(i + 1), 60, "ink", SANS_BLACK, align="c", valign="m")
        text(s, x + 312, y + ch / 2 - 96, cw - 340, 90, it.get("head", ""), 52, "ink", SANS_BLACK, valign="m")
        text(s, x + 204, y + ch / 2 + 6, cw - 240, ch / 2 - 20, ["・" + b for b in it.get("bullets", [])[:2]], 27, "sub", SANS, bold=True, line=1.5)
    pageno(s, page, total)


def s_kpi(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    rows = sl.get("rows", [])[:4]
    avail = 930 - top - 32 if sl.get("conclusion") else 1080 - top - 40
    rh = avail / max(len(rows), 1)
    for i, r in enumerate(rows):
        y = top + 16 + i * rh
        icon(s, r.get("icon", "chart"), 72, y + (rh - 96) / 2, 96, "teal")
        color = "brand-yellow" if r.get("focus") else "teal-700"
        tb = text(s, 200, y, 740, rh, r.get("value", ""), 184, color, SANS_BLACK, valign="m", align="r")
        run = tb.text_frame.paragraphs[0].add_run(); run.text = " " + r.get("unit", ""); _set_font(run, SANS_BLACK, 88, color)
        text(s, 960, y + rh / 2 - 60, 880, 60, r.get("label", ""), 44, "ink", SANS, bold=True, valign="b")
        if r.get("note"):
            text(s, 960, y + rh / 2 + 4, 880, 50, r["note"], 33, "sub", SANS, bold=True)
        if i < len(rows) - 1:
            rect(s, 72, y + rh - 1, 1776, 1, fill="line")
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_statement(prs, sl, o, page, total):
    s = blank(prs)
    circle(s, 72, 52, 56, "teal-700"); icon(s, "sprout", 72 + 10, 62, 36, "white")
    text(s, 148, 44, 900, 72, sl.get("tag", ""), 64, "ink", SANS_BLACK, valign="m")
    rect(s, 148, 122, 660, 4, fill="brand-teal")
    q = sl.get("quote", "")
    has_img = bool(sl.get("image"))
    qw = 1600 - 600 if has_img else 1600
    tb = text(s, 160, 250, qw, 420, "", 170, "ink", SANS_BLACK, align="c", valign="m", line=1.25)
    tf = tb.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = "「"; _set_font(r, SANS_BLACK, 170, "brand-yellow")
    lines = q.replace("\\n", "\n").split("\n")
    for i, ln in enumerate(lines):
        if i:
            p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.25
        r = p.add_run(); r.text = ln; _set_font(r, SANS_BLACK, 170, "ink")
    r = p.add_run(); r.text = "」"; _set_font(r, SANS_BLACK, 170, "brand-yellow")
    circle(s, 360, 740, 72, "teal-700"); icon(s, "bulb", 372, 752, 48, "white")
    text(s, 460, 720, 760 if has_img else 1300, 160, sl.get("note", ""), 48 if has_img else 58, "ink", SANS, bold=True, valign="m", line=1.5)
    picture(s, LOGO_COLOR, 1690, 950, 150, 107)
    if has_img:
        picture(s, sl["image"], 1340, 300, 480, 480)


def _panel(s, x, y, w, h, head, win):
    rect(s, x, y, w, h, fill="paper" if win else "bg", line="teal-700" if win else None, radius=8)
    rect(s, x, y, w, 80, fill="teal-700" if win else "sub", radius=8)
    rect(s, x, y + 40, w, 40, fill="teal-700" if win else "sub")
    text(s, x, y, w, 80, head, 52, "paper", SANS_BLACK, align="c", valign="m")


def _flow(s, x, y, w, nodes, sep, win):
    n = len(nodes)
    if not n:
        return
    sepw = 60; nw = (w - sepw * (n - 1)) / n
    for i, nd in enumerate(nodes):
        nx = x + i * (nw + sepw)
        rect(s, nx, y, nw, 200, fill="teal-50" if win else "line", radius=8)
        icon(s, nd.get("icon", "doc"), nx + nw / 2 - 32, y + 16, 64, "teal" if win else "ink")
        text(s, nx + 8, y + 90, nw - 16, 100, nd.get("text", ""), 28, "teal-700" if win else "ink", SANS, bold=True, align="c", valign="m")
        if i < n - 1:
            text(s, nx + nw, y, sepw, 200, sep, 44, "brand-yellow" if win else "sub", SANS_BLACK, align="c", valign="m")


def s_compare(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    y = top + 32; h = 930 - y - 24 if sl.get("conclusion") else 1040 - y
    for i, (key, win) in enumerate((("left", False), ("right", True))):
        p = sl[key]; x = 72 + i * 912; w = 864
        _panel(s, x, y, w, h, p.get("head", ""), win)
        cy = y + 110
        tw = w - 72
        if p.get("image"):
            picture(s, p["image"], x + w - 36 - 220, y + 104, 220, 220); tw -= 250
        ch = nlines(p.get("claim", ""), 46, tw) * 46 * 1.5 + 8
        text(s, x + 36, cy, tw, ch, p.get("claim", ""), 46, "ink", SANS_BLACK, line=1.35, em_color="yellow-700")
        cy += ch + 12
        if p.get("body"):
            bh = nlines(p["body"], 28, tw) * 28 * 1.6 + 8
            text(s, x + 36, cy, tw, bh, p["body"], 28, "ink", SANS, bold=True, line=1.5); cy += bh
        _flow(s, x + 36, y + h - 240, w - 72, p.get("flow", []), p.get("sep", "▶" if win else "×"), win)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_table(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    head = sl.get("head", ["御社の手作業", "MARKELINEの打ち手", "効果"])
    rows = sl.get("rows", [])[:4]
    x0 = 72; W = 1776; cw = [W * 0.34, W * 0.33, W * 0.33]
    y = top + 40
    rect(s, x0, y, W, 72, fill="teal-700", radius=8); rect(s, x0, y + 36, W, 36, fill="teal-700")
    cx = x0
    for hd, w in zip(head, cw):
        text(s, cx, y, w, 72, hd, 40, "paper", SANS, bold=True, align="c", valign="m"); cx += w
    y += 72
    rh = min(170, (1040 - y) / max(len(rows), 1))
    for i, r in enumerate(rows):
        if i % 2:
            rect(s, x0, y, W, rh, fill="bg")
        rect(s, x0, y + rh - 1, W, 1, fill="line")
        icon(s, r.get("icon", "doc"), x0 + 20, y + rh / 2 - 28, 56, "teal")
        if r.get("sub"):
            text(s, x0 + 96, y + 8, cw[0] - 110, rh / 2 - 8, r.get("task", ""), 33, "ink", SANS, bold=True, valign="b", line=1.2)
            text(s, x0 + 96, y + rh / 2 + 4, cw[0] - 110, rh / 2 - 8, r["sub"], 24, "sub", SANS, bold=True)
        else:
            text(s, x0 + 96, y + 8, cw[0] - 110, rh - 16, r.get("task", ""), 33, "ink", SANS, bold=True, valign="m", line=1.2)
        text(s, x0 + cw[0] + 20, y, 40, rh, "▶", 30, "brand-yellow", SANS_BLACK, valign="m")
        text(s, x0 + cw[0] + 60, y, cw[1] - 80, rh, r.get("action", ""), 33, "ink", SANS, bold=True, valign="m", line=1.3)
        ex = x0 + cw[0] + cw[1]
        circle(s, ex + 20, y + rh / 2 - 24, 48, "teal-700"); icon(s, "check", ex + 30, y + rh / 2 - 14, 28, "white")
        text(s, ex + 84, y, cw[2] - 100, rh, r.get("effect", ""), 38, "yellow-700", SANS_BLACK, valign="m", line=1.3)
        y += rh
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_before_after(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    y = top + 32; h = 420
    for i, (key, win) in enumerate((("before", False), ("after", True))):
        p = sl[key]; x = 72 + i * 948; w = 828
        _panel(s, x, y, w, h, p.get("head", ""), win)
        cy = y + 100
        tw = w - 72
        if p.get("image"):
            picture(s, p["image"], x + w - 36 - 220, y + 104, 220, 220); tw -= 250
        text(s, x + 36, cy, tw, 60, p.get("claim", ""), 36, "teal-700" if win else "ink", SANS_BLACK, valign="m")
        text(s, x + 60, cy + 70, tw - 24, h - (cy - y) - 80, ["・" + t for t in p.get("items", [])[:3]], 28, "ink", SANS, bold=True, line=1.5)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, px(920), px(y + h / 2 - 50), px(80), px(100))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = col("brand-yellow"); arrow.line.fill.background(); arrow.shadow.inherit = False
    fy = y + h + 24; fh = 930 - fy - 24 if sl.get("conclusion") else 1040 - fy
    fx = sl.get("effects", [])[:3]
    if fx:
        fw = (1776 - 24 * (len(fx) - 1)) / len(fx)
        for i, e in enumerate(fx):
            x = 72 + i * (fw + 24)
            rect(s, x, fy, fw, fh, fill="paper", line="teal-700", radius=8)
            circle(s, x + 24, fy + fh / 2 - 36, 72, "teal-700"); icon(s, e.get("icon", "check"), x + 40, fy + fh / 2 - 20, 40, "white")
            text(s, x + 116, fy + 12, fw - 130, fh / 2 - 8, e.get("head", ""), 44, "teal-700", SANS_BLACK, valign="b")
            text(s, x + 116, fy + fh / 2 + 4, fw - 130, fh / 2 - 12, e.get("text", ""), 29, "ink", SANS, bold=True, line=1.35)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_steps(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    steps = sl.get("steps", [])[:4]
    base = 900 if sl.get("conclusion") else 1010
    fills = ["teal-50", "teal-100", "teal-200", "brand-teal"]
    heights = [380, 460, 540, 620]
    n = len(steps); gap = 20; w = (1776 - gap * (n - 1)) / n
    # rising arrow behind the boxes
    ln = s.shapes.add_connector(1, px(72), px(base - 10), px(1848), px(top + 60))
    ln.line.color.rgb = col("brand-yellow"); ln.line.width = px(14)
    tail = ln.line._get_or_add_ln().makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "lg", "len": "lg"}); ln.line._get_or_add_ln().append(tail)
    for i, st in enumerate(steps):
        x = 72 + i * (w + gap); h = min(heights[i], base - top - 60); y = base - h
        last = i == n - 1
        rect(s, x, y, w, h, fill=fills[i], radius=8)
        circle(s, x + 26, y + 26, 56, "paper"); icon(s, st.get("icon", "run"), x + 36, y + 36, 36, "teal700")
        text(s, x + 96, y + 26, w - 120, 56, st.get("label", f"STEP{i}"), 36, "paper" if last else "teal-700", SANS_BLACK, valign="m")
        head = st.get("head", "")
        cpl = max(1, int((w - 52) / (41 * 0.97)))
        head_h = int(-(-len(plain(head)) // cpl) * 41 * 1.3) + 8
        text(s, x + 26, y + 100, w - 52, head_h, head, 41, "paper" if last else "ink", SANS_BLACK, line=1.3)
        ty = y + 100 + head_h + 12
        text(s, x + 26, ty, w - 52, max(40, y + h - 16 - ty), st.get("text", ""), 29, "paper" if last else "sub", SANS, bold=True, line=1.45)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_why_now(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    items = sl.get("items", [])[:3]
    avail = (930 if sl.get("conclusion") else 1040) - top - 40
    rh = min(260, (avail - 20 * (len(items) - 1)) / max(len(items), 1))
    y0 = top + 20 + (avail - (rh * len(items) + 20 * (len(items) - 1))) / 2
    for i, it in enumerate(items):
        y = y0 + i * (rh + 20)
        rect(s, 72, y, 1776, rh, fill="bg", radius=8)
        text(s, 100, y, 180, rh, str(i + 1), 230, "teal-700", SANS_BLACK, align="c", valign="m", wrap=False)
        circle(s, 300, y + rh / 2 - 60, 120, "teal-50"); icon(s, it.get("icon", "bulb"), 328, y + rh / 2 - 32, 64, "teal")
        text(s, 470, y + 24, 1340, 90, it.get("head", ""), 64, "ink", SANS_BLACK, valign="m", em_color="yellow-700")
        text(s, 470, y + 124, 1340, rh - 140, it.get("text", ""), 30, "sub", SANS, bold=True, line=1.5)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def _pcard(s, x, y, w, h, c, fill, fg, big=176):
    rect(s, x, y, w, h, fill=fill, radius=16)
    text(s, x + 44, y + 36, w - 88, 60, c.get("label", ""), 40, fg, SANS_BLACK, valign="m")
    tb = text(s, x + 44, y + 110, w - 88, big + 30, c.get("value", ""), big, fg, SANS_BLACK, valign="m", wrap=False)
    r = tb.text_frame.paragraphs[0].add_run(); r.text = c.get("unit", ""); _set_font(r, SANS_BLACK, 56, fg)
    text(s, x + 44, y + 110 + big + 40, w - 88, h - (110 + big + 60), c.get("text", ""), 33, fg, SANS, bold=True, line=1.4)


def s_pricing(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    y = top + 32; h = 470
    _pcard(s, 72, y, 864, h, sl["entry"], "teal-700", "paper")
    if sl.get("expand"):
        _pcard(s, 984, y, 864, h, sl["expand"], "brand-yellow", "ink")
    cmp = sl.get("compare")
    if cmp:
        cy = y + h + 32; ch = (930 if sl.get("conclusion") else 1040) - cy - 16
        rect(s, 72, cy, 1776, ch, fill="bg", radius=16)
        text(s, 116, cy + 24, 900, 50, cmp.get("label", ""), 40, "ink", SANS_BLACK, valign="m")
        tb = text(s, 116, cy + 80, 900, 130, cmp.get("value", ""), 110, "teal-700", SANS_BLACK, valign="m", wrap=False)
        r = tb.text_frame.paragraphs[0].add_run(); r.text = cmp.get("unit", ""); _set_font(r, SANS_BLACK, 56, "teal-700")
        note = ("\n" + cmp["note"]) if cmp.get("note") else ""
        text(s, 1000, cy + 24, 800, ch - 48, cmp.get("text", "") + note, 33, "ink", SANS, bold=True, valign="m", line=1.5)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_case_kpi(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    y = top + 32
    rect(s, 72, y, 700, 90, fill="teal-700", radius=8)
    rect(s, 92, y + 16, 110, 58, fill="paper", radius=4)
    text(s, 92, y + 16, 110, 58, "事例", 32, "teal-700", SANS_BLACK, align="c", valign="m")
    text(s, 220, y, 540, 90, sl.get("case", ""), 66, "paper", SANS_BLACK, valign="m")
    rows = sl.get("rows", [])[:3]
    ry = y + 120; bottom = 930 - 24 if sl.get("conclusion") else 1040
    rh = (bottom - ry) / max(len(rows), 1)
    for i, r in enumerate(rows):
        yy = ry + i * rh
        circle(s, 72, yy + rh / 2 - 48, 96, "teal-50"); icon(s, r.get("icon", "check"), 94, yy + rh / 2 - 26, 52, "teal")
        has_nums = bool(r.get("after"))
        text(s, 200, yy + (10 if has_nums else 0), 1120, 60 if has_nums else rh, r.get("label", ""), 40, "ink", SANS, bold=True, valign="t" if has_nums else "m")
        if has_nums:
            tb = text(s, 200, yy + 64, 1120, rh - 80, "", 46, "ink", SANS_BLACK, line=1.05)
            p = tb.text_frame.paragraphs[0]
            def run(t, size, color):
                rr = p.add_run(); rr.text = t; _set_font(rr, SANS_BLACK, size, color)
            if r.get("prefix"): run(r["prefix"], 46, "ink")
            if r.get("before"):
                run(r["before"], 132, "teal-700"); run(r.get("before_unit", ""), 46, "ink"); run(" ▶ ", 60, "brand-yellow")
            run(r["after"], 160, "teal-700"); run(r.get("after_unit", ""), 46, "ink")
        if i < len(rows) - 1:
            rect(s, 72, yy + rh - 1, 1250, 1, fill="line")
    side = sl.get("side", {})
    sx, sy, sw, sh = 1360, ry, 488, bottom - ry
    rect(s, sx, sy, sw, sh, fill="teal-50", radius=16)
    if side.get("image"):
        d = min(360, sw - 48, sh - 220)
        picture(s, side["image"], sx + (sw - d) / 2, sy + (sh - d - 160) / 2, d, d)
        ty = sy + (sh - d - 160) / 2 + d + 20
    else:
        circle(s, sx + sw / 2 - 100, sy + sh / 2 - 190, 200, "brand-yellow"); icon(s, side.get("icon", "target"), sx + sw / 2 - 55, sy + sh / 2 - 145, 110, "ink")
        ty = sy + sh / 2 + 40
    text(s, sx + 24, ty, sw - 48, 140, side.get("text", ""), 34, "ink", SANS_BLACK, align="c", valign="m", line=1.4)
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_case_flow(prs, sl, o, page, total):
    s = blank(prs); top = titlebar(s, sl)
    y = top + 32; h = 930 - y - 24 if sl.get("conclusion") else 1040 - y
    b = sl["before"]; a = sl["after"]
    _panel(s, 72, y, 864, h, b.get("head", ""), False)
    cy = y + 100
    tw = 792
    if b.get("image"):
        picture(s, b["image"], 72 + 864 - 36 - 220, y + 104, 220, 220); tw -= 250
    text(s, 108, cy, tw, 100, b.get("claim", ""), 36, "ink", SANS_BLACK, valign="m", line=1.3)
    cy += 110
    for it in b.get("items", [])[:3]:
        circle(s, 108, cy + 8, 36, "sub"); text(s, 108, cy + 8, 36, 36, "×", 22, "paper", SANS_BLACK, align="c", valign="m")
        text(s, 160, cy, tw - 52, 90, it, 27, "ink", SANS, bold=True, line=1.45); cy += 100
    _panel(s, 984, y, 864, h, a.get("head", ""), True)
    text(s, 1020, y + 100, 792, 50, a.get("flow_title", "一気通貫フロー図"), 36, "teal-700", SANS_BLACK, align="c", valign="m")
    rect(s, 1020, y + 152, 792, 3, fill="teal-700")
    nodes = a.get("flow", [])[:7]; n = len(nodes)
    if n:
        nw = 792 / n
        for i, nd in enumerate(nodes):
            nx = 1020 + i * nw
            circle(s, nx + nw / 2 - 42, y + 176, 84, "paper", line="brand-teal"); icon(s, nd.get("icon", "doc"), nx + nw / 2 - 22, y + 196, 44, "teal")
            text(s, nx, y + 268, nw, 40, nd.get("text", ""), 25, "teal-700", SANS, bold=True, align="c")
            if i < n - 1:
                text(s, nx + nw - 12, y + 176, 24, 84, "▶", 22, "brand-yellow", SANS_BLACK, align="c", valign="m")
    cy = y + 340
    for c in a.get("checks", [])[:3]:
        circle(s, 1020, cy + 6, 36, "teal-700"); icon(s, "check", 1028, cy + 14, 20, "white")
        text(s, 1072, cy, 740, 80, c, 27, "ink", SANS, bold=True, line=1.45); cy += 90
    conclusion(s, sl.get("conclusion"))
    if not sl.get("conclusion"):
        pageno(s, page, total)


def s_closing(prs, sl, o, page, total):
    s = blank(prs)
    bgs = rect(s, 0, 0, 1920, 1080, fill="teal-900")
    bgs.fill.gradient(); bgs.fill.gradient_angle = 160
    st = bgs.fill.gradient_stops
    st[0].color.rgb = col("teal-900"); st[0].position = 0
    st[1].color.rgb = col("teal-700"); st[1].position = 1
    text(s, 96, 64, 1200, 130, sl.get("title", "まとめ／次の一歩"), 116, "paper", SANS_BLACK, valign="m")
    picture(s, LOGO_WHITE, 1660, 60, 164, 117)
    rect(s, 96, 240, 830, 560, line="teal-400", radius=8)
    text(s, 136, 272, 750, 70, "【組むメリット】", 56, "teal-100", SANS_BLACK, valign="m")
    cy = 360
    for m in sl.get("merits", [])[:4]:
        circle(s, 136, cy + 8, 52, "brand-yellow"); icon(s, "check", 148, cy + 20, 28, "ink")
        text(s, 208, cy, 690, 70, m, 42, "paper", SANS, bold=True, valign="m"); cy += 96
    rect(s, 994, 240, 830, 560, fill="brand-yellow", radius=8)
    circle(s, 1034, 272, 60, "ink"); icon(s, "run", 1046, 284, 36, "yellow")
    text(s, 1110, 272, 680, 60, sl.get("next", {}).get("head", "【次の一歩】"), 44, "ink", SANS_BLACK, valign="m")
    rect(s, 1034, 352, 750, 2, fill="ink")
    text(s, 1034, 372, 750, 400, sl.get("next", {}).get("text", ""), 58, "ink", SANS_BLACK, line=1.4)
    msg = sl.get("message", ""); suf = sl.get("message_suffix", "を、御社に。")
    size = fit(msg + "  " + suf[:4], 118, 14, 76, 1700)
    tb = text(s, 60, 860, 1800, 160, "", size, "paper", SERIF_BLACK, align="c", valign="m", wrap=False)
    p = tb.text_frame.paragraphs[0]
    for t, c, sz in (("「", "brand-yellow", size), (msg, "paper", size), ("」", "brand-yellow", size), (suf, "paper", int(size * 0.72))):
        r = p.add_run(); r.text = t; _set_font(r, SERIF_BLACK, sz, c)


BUILDERS = {"cover": s_cover, "goals": s_goals, "issues": s_issues, "kpi": s_kpi, "statement": s_statement,
            "compare": s_compare, "table": s_table, "before_after": s_before_after, "steps": s_steps,
            "why_now": s_why_now, "pricing": s_pricing, "case_kpi": s_case_kpi, "case_flow": s_case_flow,
            "closing": s_closing}


def build(outline: dict, out: Path, base: Path | None = None) -> Path:
    global BASE_DIR
    if base is not None:
        BASE_DIR = Path(base)
    prs = Presentation()
    prs.slide_width = px(1920); prs.slide_height = px(1080)
    slides = outline["slides"]; total = len(slides)
    for i, sl in enumerate(slides, 1):
        fn = BUILDERS.get(sl.get("type"))
        if fn:
            fn(prs, sl, outline, i, total)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    _retheme_fonts(out)
    return out


def _retheme_fonts(path: Path) -> None:
    """python-pptx's default template carries Arial/Calibri in the theme and master.
    Point them at the brand font so lint_tokens passes and new text boxes inherit it."""
    import zipfile
    import shutil
    import tempfile
    tmp = Path(tempfile.mkstemp(suffix=".pptx")[1])
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/") and item.filename.endswith(".xml"):
                s = data.decode("utf-8")
                for face in ("Calibri Light", "Calibri", "Arial"):
                    s = s.replace(f'typeface="{face}"', f'typeface="{SANS}"')
                data = s.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(str(tmp), str(path))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("outline"); ap.add_argument("--out", default="deck/deck.pptx")
    ap.add_argument("--weight", choices=["bold", "black"], default="bold",
                    help="bold: 'Noto Sans JP' + bold (renders as-is in Google Slides). "
                         "black: request 'Noto Sans JP Black' faces (PowerPoint with the font installed).")
    a = ap.parse_args()
    global BLACK_FACE
    BLACK_FACE = a.weight == "black"
    outline = json.loads(Path(a.outline).read_text(encoding="utf-8"))
    out = build(outline, Path(a.out), base=Path(a.outline).resolve().parent)
    print(f"wrote {out} ({len(outline['slides'])} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
